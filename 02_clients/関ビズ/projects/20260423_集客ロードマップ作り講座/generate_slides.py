"""
関ビズ講座企画書スライド生成スクリプト
「明日からお客さんが増える！AIでつくる「うちの店の集客プラン」」
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 色定義
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x2D, 0x6A, 0x9F)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x27, 0xAE, 0x60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=WHITE):
    """スライド背景色を設定"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    """色付き矩形を追加"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    """テキストフレームに段落を追加"""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = space_before
    return p

# ========================================
# スライド1: 表紙
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 白紙
add_bg(slide, NAVY)

# アクセントライン
add_shape(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(0.08), ORANGE)

# メインタイトル
add_text_box(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(1.5),
    "明日からお客さんが増える！", 48, WHITE, True, PP_ALIGN.CENTER)

tf = add_text_box(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(1.5),
    'AIでつくる「うちの店の集客プラン」', 44, ORANGE, True, PP_ALIGN.CENTER)

# サブタイトル
add_text_box(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.8),
    "〜 無料AIツール「FELO」で、90分であなたの店の集客アクションが見つかる 〜", 22, RGBColor(0xBB, 0xCC, 0xDD), False, PP_ALIGN.CENTER)

# 日付・主催
add_text_box(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.8),
    "2026年4月23日（木）  |  関市ビジネスサポートセンター（関ビズ）", 20, RGBColor(0x88, 0x99, 0xAA), False, PP_ALIGN.CENTER)

# ========================================
# スライド2: こんなお悩みありませんか？
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.8),
    "こんなお悩み、ありませんか？", 36, WHITE, True, PP_ALIGN.CENTER)

worries = [
    ("😟", "SNSをやった方がいいのは\nわかるけど、何を投稿すれば…"),
    ("😟", "チラシを配っても\n反応がわからない…"),
    ("😟", "お客さんが来ない\n時間帯をなんとかしたい…"),
    ("😟", "集客って結局\n何から始めればいいの？"),
]

for i, (emoji, text) in enumerate(worries):
    left = Inches(0.8 + i * 3.1)
    # カード背景
    shape = add_shape(slide, left, Inches(1.8), Inches(2.8), Inches(4.5), LIGHT_BLUE)
    shape.shadow.inherit = False
    # テキスト
    add_text_box(slide, left + Inches(0.3), Inches(2.0), Inches(2.2), Inches(0.8),
        emoji, 60, DARK_GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), Inches(3.2), Inches(2.4), Inches(2.5),
        text, 22, DARK_GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.8),
    "▶ この講座で、あなたの店に合った「具体的なアクション」が見つかります", 24, ORANGE, True, PP_ALIGN.CENTER)

# ========================================
# スライド3: 講座概要
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.8),
    "講座概要", 36, WHITE, True, PP_ALIGN.CENTER)

# 左カード: 午前の部
shape = add_shape(slide, Inches(1), Inches(1.8), Inches(5.2), Inches(3.0), LIGHT_BLUE)
add_text_box(slide, Inches(1.3), Inches(1.9), Inches(4.6), Inches(0.6),
    "午前の部", 28, BLUE, True, PP_ALIGN.CENTER)
tf = add_text_box(slide, Inches(1.5), Inches(2.6), Inches(4.2), Inches(2.0),
    "⏰  10:00 〜 11:30（90分）", 22, DARK_GRAY)
add_paragraph(tf, "👥  定員 6名", 22, DARK_GRAY)

# 右カード: 午後の部
shape = add_shape(slide, Inches(7.1), Inches(1.8), Inches(5.2), Inches(3.0), LIGHT_BLUE)
add_text_box(slide, Inches(7.4), Inches(1.9), Inches(4.6), Inches(0.6),
    "午後の部", 28, BLUE, True, PP_ALIGN.CENTER)
tf = add_text_box(slide, Inches(7.6), Inches(2.6), Inches(4.2), Inches(2.0),
    "⏰  13:00 〜 14:30（90分）", 22, DARK_GRAY)
add_paragraph(tf, "👥  定員 6名", 22, DARK_GRAY)

# 下部情報
info_items = [
    "📱  持ち物：スマホまたはPC（各自の通信回線を使用）",
    "🤖  使用ツール：FELO（無料AIツール・アカウント登録不要）",
    "📝  成果物：1枚アクションシート（ネット施策＋リアル施策）",
    "💡  事前準備：今やっている集客方法を整理しておくだけでOK",
]

tf = add_text_box(slide, Inches(1.2), Inches(5.2), Inches(11), Inches(2.2),
    info_items[0], 20, DARK_GRAY)
for item in info_items[1:]:
    add_paragraph(tf, item, 20, DARK_GRAY, space_before=Pt(10))

# ========================================
# スライド4: この講座で得られること
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.8),
    "90分で得られること", 36, WHITE, True, PP_ALIGN.CENTER)

outcomes = [
    ("1", "お客様の行動が\n見えるようになる",
     "AIがあなたのお客様の\n1日の行動をシミュレーション。\n「いつ」「どこで」情報を\n届ければいいかがわかります。"),
    ("2", "今のSNS・集客の\n改善点がわかる",
     "今やっているInstagramや\nチラシの効果をAIが診断。\n「何を変えればいいか」が\n具体的にわかります。"),
    ("3", "明日からの\nアクションが決まる",
     "ネット（SNS）とリアル\n（店頭・口コミ）に分けて、\nコスト0円で今すぐできる\nアクションを持ち帰れます。"),
]

for i, (num, title, desc) in enumerate(outcomes):
    left = Inches(0.8 + i * 4.1)
    # 番号丸
    shape = add_shape(slide, left + Inches(1.2), Inches(1.5), Inches(0.8), Inches(0.8), ORANGE)
    add_text_box(slide, left + Inches(1.2), Inches(1.5), Inches(0.8), Inches(0.8),
        num, 32, WHITE, True, PP_ALIGN.CENTER)
    # タイトル
    add_text_box(slide, left, Inches(2.5), Inches(3.6), Inches(1.2),
        title, 24, NAVY, True, PP_ALIGN.CENTER)
    # 説明
    add_text_box(slide, left + Inches(0.2), Inches(3.9), Inches(3.2), Inches(3.0),
        desc, 18, DARK_GRAY, False, PP_ALIGN.CENTER)

# ========================================
# スライド5: FELOとは？
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.8),
    "使うのはコレだけ！ 無料AIツール「FELO」", 36, WHITE, True, PP_ALIGN.CENTER)

# 左側：FELOの説明
shape = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0), LIGHT_BLUE)
tf = add_text_box(slide, Inches(1.2), Inches(2.0), Inches(5.0), Inches(4.5),
    "FELOってなに？", 28, BLUE, True)
add_paragraph(tf, "", 12, DARK_GRAY)
add_paragraph(tf, "✅  質問を入力するだけで、AIが最新情報を", 20, DARK_GRAY)
add_paragraph(tf, "    調べてまとめてくれる無料ツール", 20, DARK_GRAY)
add_paragraph(tf, "", 12, DARK_GRAY)
add_paragraph(tf, "✅  アカウント登録なしですぐ使える", 20, DARK_GRAY)
add_paragraph(tf, "", 12, DARK_GRAY)
add_paragraph(tf, "✅  スマホのブラウザから開くだけ", 20, DARK_GRAY)
add_paragraph(tf, "", 12, DARK_GRAY)
add_paragraph(tf, "✅  講座では講師と一緒に操作するので", 20, DARK_GRAY)
add_paragraph(tf, "    初めてでも安心", 20, DARK_GRAY)

# 右側：使い方
shape = add_shape(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(5.0), RGBColor(0xFF, 0xF3, 0xE0))
tf = add_text_box(slide, Inches(7.6), Inches(2.0), Inches(4.5), Inches(4.5),
    "講座での使い方", 28, ORANGE, True)
add_paragraph(tf, "", 12, DARK_GRAY)
add_paragraph(tf, "① 講師がスクリーンで実演", 22, DARK_GRAY, True)
add_paragraph(tf, "    ↓", 22, DARK_GRAY)
add_paragraph(tf, "② 参加者も一緒に同じ操作", 22, DARK_GRAY, True)
add_paragraph(tf, "    ↓", 22, DARK_GRAY)
add_paragraph(tf, "③ AIの回答をシートにメモ", 22, DARK_GRAY, True)
add_paragraph(tf, "", 14, DARK_GRAY)
add_paragraph(tf, "※ 難しい操作は一切ありません", 18, ORANGE)

# ========================================
# スライド6: タイムテーブル
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.8),
    "90分の流れ", 36, WHITE, True, PP_ALIGN.CENTER)

schedule = [
    ("0:00", "5分", "オープニング", "今日のゴール説明", LIGHT_GRAY),
    ("0:05", "10分", "ワーク①", "自店の現状を書き出す\n（SNS・集客の棚卸し）", LIGHT_BLUE),
    ("0:15", "30分", "ワーク②", "講師と一緒にFELOで分析\n（3ステップ）", RGBColor(0xFF, 0xF3, 0xE0)),
    ("0:45", "20分", "ワーク③", "アクションシート完成\n（ネット/リアルに分けて）", LIGHT_BLUE),
    ("1:05", "10分", "グループシェア", "全員でシートを共有\n互いの施策から学び合う", RGBColor(0xE8, 0xF8, 0xE8)),
    ("1:15", "10分", "宣言タイム", "「明日やること1つ」\nを全員が発表", RGBColor(0xFF, 0xF3, 0xE0)),
    ("1:25", "5分", "まとめ", "関ビズ個別相談の案内", LIGHT_GRAY),
]

for i, (time, duration, title, desc, color) in enumerate(schedule):
    top = Inches(1.5 + i * 0.82)
    # 背景帯
    add_shape(slide, Inches(0.8), top, Inches(11.7), Inches(0.75), color)
    # 時間
    add_text_box(slide, Inches(0.9), top + Inches(0.1), Inches(1.0), Inches(0.55),
        time, 16, DARK_GRAY, False, PP_ALIGN.CENTER)
    # 所要時間
    add_text_box(slide, Inches(1.9), top + Inches(0.1), Inches(1.0), Inches(0.55),
        duration, 16, BLUE, True, PP_ALIGN.CENTER)
    # タイトル
    add_text_box(slide, Inches(3.0), top + Inches(0.1), Inches(2.8), Inches(0.55),
        title, 18, NAVY, True)
    # 説明
    add_text_box(slide, Inches(6.0), top + Inches(0.05), Inches(6.2), Inches(0.65),
        desc, 15, DARK_GRAY)

# ========================================
# スライド7: 3ステップの詳細
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.8),
    "AIに聞く3つの質問（ワーク②の詳細）", 36, WHITE, True, PP_ALIGN.CENTER)

steps = [
    ("STEP 1", "お客様の「1日」を\n見える化する",
     "理想のお客様がいつスマホを\n見るか、いつお店の情報に\n触れるかをAIが分析\n\n→ 今のInstagram投稿の\n   診断もしてくれる",
     "8分", BLUE),
    ("STEP 2", "「知る→気になる→\n 行ってみる」を分析",
     "お客様が店を知ってから\n来店するまでの心の動きと、\n最適な伝え方をAIが提案\n\n→ 背中を押す「最後の一言」\n   も教えてくれる",
     "10分", GREEN),
    ("STEP 3", "明日からできる\nアクションプラン",
     "ネット施策とリアル施策に\n分けて、コスト0円で\nできるアクションをAIが提案\n\n→ 明日/1週間/1ヶ月の\n   3段階で具体的に",
     "12分", ORANGE),
]

for i, (step, title, desc, time, color) in enumerate(steps):
    left = Inches(0.6 + i * 4.2)
    # カード
    shape = add_shape(slide, left, Inches(1.5), Inches(3.9), Inches(5.5), LIGHT_GRAY)
    # ステップラベル
    shape = add_shape(slide, left, Inches(1.5), Inches(3.9), Inches(0.7), color)
    add_text_box(slide, left, Inches(1.5), Inches(2.5), Inches(0.7),
        f"  {step}", 22, WHITE, True)
    add_text_box(slide, left + Inches(2.5), Inches(1.5), Inches(1.3), Inches(0.7),
        time, 18, WHITE, False, PP_ALIGN.RIGHT)
    # タイトル
    add_text_box(slide, left + Inches(0.3), Inches(2.4), Inches(3.3), Inches(1.2),
        title, 22, NAVY, True, PP_ALIGN.CENTER)
    # 説明
    add_text_box(slide, left + Inches(0.3), Inches(3.7), Inches(3.3), Inches(3.0),
        desc, 16, DARK_GRAY)

# ========================================
# スライド8: グループシェア
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.8),
    "参加者同士で学び合う「グループシェア」", 36, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
    "6人全員でアクションシートを共有し、互いの集客アイデアからヒントをもらう", 22, DARK_GRAY, False, PP_ALIGN.CENTER)

share_steps = [
    ("1", "自分のシートを共有",
     "「うちの店は今〇〇をやっていて、\nAIに聞いたら〇〇がわかりました」\n→ 1人1分で発表"),
    ("2", "お客さん目線で\nフィードバック",
     "「その投稿、○○時に見たら\n行きたくなります！」\n→ 異業種の視点が新鮮な気づきに"),
    ("3", "他の人のアイデアを\n自分のシートに追記",
     "「これ自分もやりたい！」\nと思ったものは遠慮なくメモ\n→ 6人分のアイデアが全員の財産に"),
]

for i, (num, title, desc) in enumerate(share_steps):
    left = Inches(0.8 + i * 4.1)
    shape = add_shape(slide, left, Inches(2.6), Inches(3.8), Inches(4.2), RGBColor(0xE8, 0xF8, 0xE8))
    # 番号
    shape = add_shape(slide, left + Inches(1.3), Inches(2.8), Inches(0.8), Inches(0.8), GREEN)
    add_text_box(slide, left + Inches(1.3), Inches(2.8), Inches(0.8), Inches(0.8),
        num, 32, WHITE, True, PP_ALIGN.CENTER)
    # タイトル
    add_text_box(slide, left + Inches(0.2), Inches(3.8), Inches(3.4), Inches(1.0),
        title, 20, NAVY, True, PP_ALIGN.CENTER)
    # 説明
    add_text_box(slide, left + Inches(0.3), Inches(4.9), Inches(3.2), Inches(1.8),
        desc, 16, DARK_GRAY)

# ========================================
# スライド9: 参加者が持ち帰るもの
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.8),
    "参加者が持ち帰るもの", 36, WHITE, True, PP_ALIGN.CENTER)

# メイン成果物
shape = add_shape(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(2.5), RGBColor(0xFF, 0xF3, 0xE0))
add_text_box(slide, Inches(1.5), Inches(1.7), Inches(10.3), Inches(0.6),
    "📝  1枚アクションシート", 28, ORANGE, True)
tf = add_text_box(slide, Inches(1.5), Inches(2.4), Inches(10.3), Inches(1.5),
    "・自店の現状分析（今やっている集客の棚卸し）", 20, DARK_GRAY)
add_paragraph(tf, "・ネット施策 3つ（明日 / 1週間以内 / 1ヶ月以内）", 20, DARK_GRAY)
add_paragraph(tf, "・リアル施策 3つ（明日 / 1週間以内 / 1ヶ月以内）", 20, DARK_GRAY)
add_paragraph(tf, "・「明日やること1つ」の宣言", 20, DARK_GRAY)

# サブ成果物
items = [
    ("🤖  FELOの使い方", "講座後も自宅で使える。\n季節やメニューが変われば\n新しいプランが作れる"),
    ("💡  仲間のアイデア", "6人分の集客アイデアを\nグループシェアで共有。\n異業種の視点がヒントに"),
    ("🔗  関ビズの個別相談", "講座後もサポート。\nアクション実行の相談や\n次のステップを一緒に"),
]

for i, (title, desc) in enumerate(items):
    left = Inches(0.8 + i * 4.1)
    shape = add_shape(slide, left, Inches(4.5), Inches(3.8), Inches(2.5), LIGHT_BLUE)
    add_text_box(slide, left + Inches(0.3), Inches(4.7), Inches(3.2), Inches(0.6),
        title, 20, BLUE, True)
    add_text_box(slide, left + Inches(0.3), Inches(5.3), Inches(3.2), Inches(1.5),
        desc, 16, DARK_GRAY)

# ========================================
# スライド10: 講座後のフォロー
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.8),
    "講座後のフォロー体制", 36, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
    "講座は「始まり」。実行と改善を関ビズが継続サポート", 22, DARK_GRAY, False, PP_ALIGN.CENTER)

follow_items = [
    ("講座当日", "希望者にその場で\n関ビズ個別相談を予約", "→ すぐに次のステップへ"),
    ("1週間後", "参加者へフォロー連絡\nアクション進捗の確認", "→ 実行のモチベーション維持"),
    ("1ヶ月後", "成果ヒアリング\n（事例化の許可取り）", "→ 次回講座の実績にも活用"),
]

for i, (timing, action, effect) in enumerate(follow_items):
    left = Inches(0.8 + i * 4.1)
    shape = add_shape(slide, left, Inches(2.6), Inches(3.8), Inches(4.0), LIGHT_GRAY)
    # タイミング
    shape = add_shape(slide, left + Inches(0.8), Inches(2.8), Inches(2.2), Inches(0.7), BLUE)
    add_text_box(slide, left + Inches(0.8), Inches(2.8), Inches(2.2), Inches(0.7),
        timing, 22, WHITE, True, PP_ALIGN.CENTER)
    # アクション
    add_text_box(slide, left + Inches(0.3), Inches(3.8), Inches(3.2), Inches(1.5),
        action, 18, DARK_GRAY, False, PP_ALIGN.CENTER)
    # 効果
    add_text_box(slide, left + Inches(0.3), Inches(5.5), Inches(3.2), Inches(0.8),
        effect, 16, GREEN, True, PP_ALIGN.CENTER)

# ========================================
# スライド11: まとめ
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(0.08), ORANGE)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(1.0),
    "明日からお客さんが増える！", 44, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(0.8),
    'AIでつくる「うちの店の集客プラン」', 40, ORANGE, True, PP_ALIGN.CENTER)

summary_items = [
    "✅  無料AIツール「FELO」で、自分の店に合った集客法を発見",
    "✅  ネット × リアル、両方のアクションプランを持ち帰り",
    "✅  参加者同士のシェアで、6人分のアイデアが全員の財産に",
    "✅  講座後も関ビズが継続サポート",
]

tf = add_text_box(slide, Inches(2.5), Inches(3.5), Inches(8.3), Inches(3.0),
    summary_items[0], 22, WHITE)
for item in summary_items[1:]:
    add_paragraph(tf, item, 22, WHITE, space_before=Pt(14))

add_text_box(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.8),
    "2026年4月23日（木）  午前10:00 / 午後13:00  |  各回定員6名", 24, RGBColor(0xBB, 0xCC, 0xDD), False, PP_ALIGN.CENTER)

# ========================================
# 保存
# ========================================
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "関ビズ企画書_集客プラン講座_0423.pptx")
prs.save(output_path)
print(f"スライド生成完了: {output_path}")
print(f"全{len(prs.slides)}枚")
